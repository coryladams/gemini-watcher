#!/usr/bin/env python3
"""
gemini_watcher.py
Watches a folder for new files, processes them via `gemini` CLI,
and routes the output directly to your Obsidian vault.

Optimized for: Gemini CLI + Obsidian + WSL
"""

import os
import sys
import time
import json
import logging
import subprocess
import pypdf
import shutil
import hashlib
import difflib
import docx
import pptx
import openpyxl
from pathlib import Path
from datetime import datetime
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

# --- Vault Configuration (Hardcoded for Cadams) -------------------------------
VAULT_ROOT = "/mnt/c/Users/cadams/Documents/Obsidian/RegenMed"
DEFAULT_WATCH_DIR = "/mnt/c/Users/cadams/Downloads/VaultDrop"
DEFAULT_ARCHIVE_DIR = "/mnt/c/Users/cadams/Downloads/VaultDrop/Archive"
DEFAULT_FAILED_DIR = "/mnt/c/Users/cadams/Downloads/VaultDrop/Failed"
ATTACHMENTS_DIR = "Research/Attachments"

CONFIG_PATH = Path(__file__).parent / "gemini_config.json"

def load_config():
    defaults = {
        "watch_dir": DEFAULT_WATCH_DIR,
        "vault_root": VAULT_ROOT,
        "archive_dir": DEFAULT_ARCHIVE_DIR,
        "failed_dir": DEFAULT_FAILED_DIR,
        "archive_originals": True,
        "smart_routing": True,
        "model": "gemini-2.5-flash",
        "routes": {
            "research": {
                "extensions": ["pdf", "txt", "docx", "pptx", "rtf", "epub"],
                "destination": "Research",
                "instructions": "Summarize this technical document. Identify key concepts, tools, and methodologies. Use callouts for summaries and warnings."
            },
            "markdown": {
                "extensions": ["md"],
                "destination": "Research",
                "instructions": "Read the provided markdown document. You MUST output the ENTIRE original markdown content exactly as it is, but insert a brief Obsidian callout summary (> [!summary]) at the beginning of each major section summarizing that section."
            },
            "code": {
                "extensions": ["py", "cpp", "h", "js", "ts", "ino", "sh", "yaml", "xml", "sql"],
                "destination": "Coding 1s & 0s/Code",
                "instructions": "Explain this source code. Detail the architecture, key functions, and any dependencies found."
            },
            "data": {
                "extensions": ["csv", "json", "xlsx"],
                "destination": "Research/Data",
                "instructions": "Analyze this dataset. Identify patterns, anomalies, and provide a high-level statistical summary."
            },
            "images": {
                "extensions": ["png", "jpg", "jpeg", "webp"],
                "destination": "Research/Images",
                "instructions": "Analyze this image. If it is a diagram or schematic, explain its components and flow. If it is a screenshot of code or text, transcribe and summarize it. Describe the visual content in detail."
            },
            "media": {
                "extensions": ["mp3", "mp4", "wav", "m4a", "mov"],
                "destination": "Research/Media",
                "instructions": "Analyze this media file. For audio, provide a transcript and summary. For video, describe the visual content, any spoken words, and summarize the key events."
            },
            "webpages": {
                "extensions": ["html", "htm", "mhtml", "mht", "xml", "url", "webloc", "desktop"],
                "destination": "Research/Webclips",
                "instructions": "Extract and summarize the main content of this webpage. Ignore navigation menus, ads, and boilerplate text. Identify the core argument, key facts, and any cited sources. (If provided a URL, fetch it first)."
            }
        }
    }
    
    if CONFIG_PATH.exists():
        with open(CONFIG_PATH) as f:
            user_config = json.load(f)
            defaults.update(user_config)
    return defaults

# --- WSL ↔ Windows Path Translation -----------------------------------------
def wsl_to_windows(path: str) -> str:
    p = Path(path).parts
    if len(p) >= 3 and p[0] == "/" and p[1] == "mnt":
        drive = p[2].upper() + ":\\"
        rest = "\\".join(p[3:])
        return drive + rest
    return path

def resolve_path(path: str) -> Path:
    if path.startswith("/"):
        return Path(path)
    elif len(path) >= 2 and path[1] == ":":
        drive = path[0].lower()
        rest = path[3:].replace("\\", "/")
        return Path(f"/mnt/{drive}/{rest}")
    return Path(path).expanduser()

# --- Vault Indexer -----------------------------------------------------------
class VaultIndexer:
    def __init__(self, vault_root: str):
        self.vault_root = resolve_path(vault_root)
        self.cache_file = Path(__file__).parent / ".vault_index.json"
        self.ignore_dirs = {".git", ".obsidian", ".trash", "_attachments"}
        self.cache_expiry = 86400  # 24 hours

    def get_map(self) -> list[str]:
        return self.get_cache().get("folders", [])
        
    def get_cache(self) -> dict:
        if not self.vault_root.exists():
            return {"folders": [], "files": {}, "tags": []}
            
        # Check cache validity
        if self.cache_file.exists():
            try:
                with open(self.cache_file, "r") as f:
                    cache = json.load(f)
                    if time.time() - cache.get("timestamp", 0) < self.cache_expiry:
                        return cache
            except Exception:
                pass
                
        # Rebuild cache
        folders = []
        files = {}
        tags = set()
        import re
        tag_pattern = re.compile(r"^  - ([a-zA-Z0-9_-]+)$")
        
        for root, dirs, filenames in os.walk(self.vault_root):
            dirs[:] = [d for d in dirs if d not in self.ignore_dirs and not d.startswith(".")]
            try:
                rel_path = Path(root).relative_to(self.vault_root)
                if str(rel_path) != ".":
                    folders.append(str(rel_path).replace("\\", "/"))
            except ValueError:
                pass
                
            for f in filenames:
                if f.endswith(".md"):
                    basename = Path(f).stem
                    filepath = Path(root) / f
                    try:
                        rel_file = str(filepath.relative_to(self.vault_root)).replace("\\", "/")
                        files[basename] = rel_file
                    except ValueError:
                        continue
                    
                    try:
                        with open(filepath, "r", encoding="utf-8", errors="ignore") as fp:
                            in_tags = False
                            for i, line in enumerate(fp):
                                if i > 40: break
                                if line.startswith("tags:"):
                                    in_tags = True
                                elif in_tags and line.startswith("  - "):
                                    match = tag_pattern.match(line.rstrip())
                                    if match: tags.add(match.group(1))
                                elif in_tags and not line.startswith(" "):
                                    in_tags = False
                    except Exception:
                        pass
                
        cache_data = {
            "timestamp": time.time(),
            "folders": folders,
            "files": files,
            "tags": list(tags)
        }
        try:
            with open(self.cache_file, "w") as f:
                json.dump(cache_data, f)
        except Exception as e:
            logging.warning(f"Could not save vault index cache: {e}")
            
        return cache_data

    def invalidate_cache(self):
        if self.cache_file.exists():
            try:
                self.cache_file.unlink()
            except Exception:
                pass

    def fuzzy_match_folder(self, target: str, folders: list[str]) -> str | None:
        if not target or target == "Inbox":
            return None
        # Exact match on full path
        if target in folders:
            return target
            
        # Exact match on folder basename (e.g. 'Boba 2.0' matching 'Boba/Boba 2.0')
        target_lower = target.lower()
        for folder in folders:
            if Path(folder).name.lower() == target_lower:
                return folder
                
        # Fuzzy match fallback
        matches = difflib.get_close_matches(target, folders, n=1, cutoff=0.6)
        return matches[0] if matches else None

# --- Read File (with PDF & URL support) -------------------------------------
def extract_docx_text(path: Path) -> str:
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_pptx_text(path: Path) -> str:
    prs = pptx.Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_xlsx_text(path: Path) -> str:
    wb = openpyxl.load_workbook(path, data_only=True)
    text = []
    for sheet in wb.worksheets:
        text.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            text.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
    return "\n".join(text)

def read_file_content(file_path: Path) -> str | None:
    try:
        suffix = file_path.suffix.lower()
        if suffix in [".pdf", ".mp3", ".mp4", ".wav", ".m4a", ".mov"]:
            # Return None to trigger raw binary file processing by Gemini CLI
            return None
        elif suffix == ".docx":
            return extract_docx_text(file_path)
        elif suffix == ".pptx":
            return extract_pptx_text(file_path)
        elif suffix == ".xlsx":
            return extract_xlsx_text(file_path)
        elif suffix in [".mhtml", ".mht"]:
            # MHTML files can be processed natively by Gemini
            return None
        elif suffix == ".url":
            content = file_path.read_text(encoding="utf-8", errors="replace")
            for line in content.splitlines():
                if line.upper().startswith("URL="):
                    url = line[4:].strip()
                    return f"Please use your web fetch capabilities to read and analyze this URL:\n{url}"
            return content
        else:
            return file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        logging.warning(f"  Could not read {file_path.name}: {e}")
        return None

# --- Clean Output -----------------------------------------------------------
def clean_gemini_output(text: str) -> str:
    """Removes CLI warnings, rogue code fences, and conversational filler that break Obsidian."""
    lines = text.splitlines()
    
    # 1. Filter out known CLI warnings
    filtered = [l for l in lines if "MCP issues detected" not in l and "Run /mcp list" not in l]
    text = "\n".join(filtered).strip()

    # 2. Remove a single leading/trailing code block if the model wrapped everything
    if text.startswith("```"):
        lines = text.splitlines()
        if lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        text = "\n".join(lines).strip()

    # 3. Fix yaml code blocks mistakenly used for frontmatter
    if text.startswith("yaml\n---"):
        text = text[5:]
    elif text.startswith("yaml\n"):
        text = "---\n" + text[5:]

    # 4. Strip conversational filler before the first frontmatter ---
    lines = text.splitlines()
    if lines and not lines[0].strip().startswith("---"):
        for i, line in enumerate(lines):
            if line.strip() == "---":
                # Found the likely start of frontmatter
                text = "\n".join(lines[i:])
                break
        
    # 5. Enforce strict Obsidian frontmatter (must start with ---)
    if not text.startswith("---"):
        if text.startswith("title:") or text.startswith("create-date:"):
            text = "---\n" + text
            lines = text.splitlines()
            for i, line in enumerate(lines):
                if line.strip() == "" or line.startswith("#") or line.startswith(">"):
                    lines.insert(i, "---")
                    text = "\n".join(lines)
                    break
        else:
            text = f"---\ntitle: \"Processed Document\"\ncreate-date: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n---\n\n" + text

    return text

# --- Process with Gemini CLI -------------------------------------------------
def process_with_gemini(file_path: Path, content: str | None, route: dict, config: dict) -> str:
    instructions = route.get("instructions", "Analyze this file.")
    
    indexer = VaultIndexer(config["vault_root"])
    cache = indexer.get_cache()
    folders = cache.get("folders", [])
    existing_tags = cache.get("tags", [])
    existing_files = list(cache.get("files", {}).keys())
    
    import re
    words = set(re.findall(r'\w+', file_path.stem.lower()))
    words = {w for w in words if len(w) > 3}
    related_files = set()
    if words:
        for ef in existing_files:
            ef_words = set(re.findall(r'\w+', ef.lower()))
            if words & ef_words:
                related_files.add(ef)
    related_files = list(related_files)[:20]
    
    related_context = ""
    if related_files:
        related_list = ", ".join(f"[[{f}]]" for f in related_files if f != file_path.stem)
        if related_list:
            related_context = f"\n   These existing vault notes share keywords with the file name. Link to them if relevant:\n   {related_list}"
            
    sibling_files = config.get("sibling_files", [])
    if sibling_files:
        siblings_list = ", ".join(f"[[{Path(s).stem}]]" for s in sibling_files if s != file_path.name)
        if siblings_list:
            related_context += f"\n   Files dropped in the same batch (Link to these in a 'Related' section):\n   {siblings_list}"

    tags_list = ", ".join(existing_tags[:100])
    tags_rule = f"Generate 3-5 highly relevant tags. PREFER reusing these existing tags: [{tags_list}]. If no existing tag fits, create a new one." if tags_list else "Generate 3-5 highly relevant tags based on the document's content and Original File Name."

    extra_instructions = f"""9. Use Obsidian callouts (> [!summary], etc.) for key sections. Ensure distinct callouts are separated by a completely empty line (no `>`).
10. If summarizing a multi-page PDF, embed each specific page alongside its summary using Obsidian syntax: `![[{file_path.name}#page=X]]` where X is the page number."""
    if config.get("smart_routing", False):
        folder_list = "\n".join(f"- {f}" for f in folders[:500])  # limit to 500
        extra_instructions = f"""9. SMART ROUTING: Choose the most appropriate existing folder for this document from the vault map below.
   Include a `destination` key in the YAML frontmatter with the chosen path (e.g., `destination: Polaris/Schematics`).
   If no specific project folder matches, use `{route.get("destination", "Research")}`.
10. Use Obsidian callouts (> [!summary], etc.) for key sections. Ensure distinct callouts are separated by a completely empty line (no `>`).
11. If summarizing a multi-page PDF, embed each specific page alongside its summary using Obsidian syntax: `![[{file_path.name}#page=X]]` where X is the page number.

Vault Map (Existing Folders):
{folder_list}
"""
    
    prompt = f"""<system_instruction>
You are an expert Research Archivist. 
Format the output as Obsidian-optimized Markdown.

CRITICAL FORMATTING RULES:
1. OUTPUT ONLY the Markdown file. DO NOT include any conversational filler (e.g. "I will begin by...").
2. You MUST start your response with exactly three dashes `---` on the very first line to begin the YAML frontmatter.
3. You MUST close the frontmatter block with exactly three dashes `---` on its own line.
4. DO NOT wrap the frontmatter in a markdown code block (no ``` or ```yaml).
5. YAML Frontmatter must use this EXACT structure (NO '#' symbols in the tags list):
---
title: "[Descriptive Title]"
create-date: {datetime.now().strftime("%Y-%m-%d %H:%M")}
type: [Choose most appropriate: reference, meeting, project, tutorial, logs]
Project: [Infer from filename/content, or 'Home']
tags:
  - tag1
  - tag2
destination: [Destination Path]
update_backlinks: []
status: complete
---
6. TAGGING RULES: {tags_rule}
7. LINKING: You MUST add a '## Related Notes' section at the bottom of the document and add wikilinks to related concepts if appropriate.{related_context}
8. BACKLINKS: If you strongly believe an existing note in the vault should be updated to point back to this new note, include its EXACT name (without the .md extension) in the `update_backlinks` list in the YAML frontmatter. Only do this for highly relevant existing notes.
{extra_instructions}
</system_instruction>

Task: {instructions}
Original File Name: {file_path.name} (Use this file name as a strong hint for choosing the correct routing destination from the Vault Map)

"""
    
    cmd_base = ["gemini"]
    
    # Define fallback models (Ordered by Capability vs Cost)
    fallback_models = [
        "gemini-2.5-flash",
        "gemini-3.1-flash-preview",
        "gemini-3.1-flash-lite-preview",
        "gemini-2.0-flash",
        "gemini-2.5-flash-lite",
        "gemini-1.5-pro"
    ]
    
    # Always put the configured model first, followed by the rest
    models_to_try = []
    if "model" in config:
        models_to_try.append(config["model"])
        
    for m in fallback_models:
        if m not in models_to_try:
            models_to_try.append(m)
    
    cmd_args = []
    # If content is None, it means it's a binary/image file.
    if content is None:
        # We copy it locally to the project dir so Gemini CLI can read it without workspace errors
        local_copy = Path(__file__).parent / file_path.name
        shutil.copy2(file_path, local_copy)
        
        cmd_args.append(f"@{local_copy.resolve()}")
        cmd_args.append(prompt)
    else:
        # For text files, we append the content to the prompt
        max_chars = 500000 
        clean_content = content.replace("\u0000", "")
        prompt += f"\n--- FILE: {file_path.name} ---\n{clean_content[:max_chars]}\n--- END ---\n"
        cmd_args.append(prompt)

    try:
        last_error = ""
        for model in models_to_try:
            retries = 2
            success = False
            for attempt in range(retries):
                cmd = cmd_base + ["--model", model] + cmd_args
                
                # Copy environment to pass down variables like API keys and inject bypass flag
                env = os.environ.copy()
                env["GEMINI_BYPASS_ROUTER"] = "1"
                
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=300,
                    encoding="utf-8",
                    cwd=Path(__file__).parent,
                    env=env
                )
                
                if result.returncode == 0:
                    if content is None:
                        local_copy.unlink(missing_ok=True)
                    return clean_gemini_output(result.stdout.strip())
                
                last_error = result.stderr.strip()
                
                # Check for known transient/network errors
                transient_errors = ["AbortError", "fetch failed", "503", "504", "timeout"]
                is_transient = any(err in last_error for err in transient_errors)
                
                if is_transient and attempt < retries - 1:
                    logging.warning(f"  ⚠️ Transient error with model '{model}': {last_error.splitlines()[-1] if last_error else 'Unknown'}. Retrying ({attempt+1}/{retries})...")
                    time.sleep(5)
                    continue
                else:
                    break # Break retry loop, move to next model or fail
                    
            # If the model is not found or we exhausted retries on transient errors, try the next model
            not_found_errors = ["ModelNotFoundError", "404", "not found"]
            if any(err in last_error.lower() for err in not_found_errors) or is_transient:
                logging.warning(f"  ⚠️ Model '{model}' failed or unavailable, trying next...")
                continue
            else:
                # If it's a different persistent error, fail fast
                if content is None:
                    local_copy.unlink(missing_ok=True)
                raise RuntimeError(f"gemini command failed with {model}: {last_error}")
                
        if content is None:
            local_copy.unlink(missing_ok=True)
        raise RuntimeError(f"All configured models failed. Last error: {last_error}")
        
    except FileNotFoundError:
        if 'local_copy' in locals():
            local_copy.unlink(missing_ok=True)
        raise RuntimeError("Gemini CLI not found. Run 'npm install -g @google/gemini-cli'.")

# --- Routing, Archiving & Embedding ------------------------------------------
def route_output(original_path: Path, output: str, route: dict, config: dict):
    target_dest = route.get("destination", "Inbox")
    
    if config.get("smart_routing", False):
        import re
        match = re.search(r"^destination:\s*(.+)$", output, re.MULTILINE | re.IGNORECASE)
        if match:
            suggested_dest = match.group(1).strip().strip("'\"")
            indexer = VaultIndexer(config["vault_root"])
            folders = indexer.get_map()
            matched_dest = indexer.fuzzy_match_folder(suggested_dest, folders)
            if matched_dest:
                target_dest = matched_dest
            else:
                logging.info(f"  ⚠️ Could not match suggested destination '{suggested_dest}', using default '{target_dest}'")

    # 1. Write the markdown note
    dest_dir = resolve_path(config["vault_root"]) / target_dest
    if not dest_dir.exists():
        dest_dir.mkdir(parents=True, exist_ok=True)
        if config.get("smart_routing", False):
            VaultIndexer(config["vault_root"]).invalidate_cache()
            logging.info(f"  🔄 Created new folder '{target_dest}', invalidated index cache")

    dest_file = dest_dir / f"{original_path.stem}.md"
    if dest_file.exists():
        timestamp = datetime.now().strftime("%H%M%S")
        dest_file = dest_dir / f"{original_path.stem}-{timestamp}.md"

    # Append the source link if it's an internet shortcut
    if original_path.suffix.lower() in [".url", ".webloc"]:
        try:
            url_content = original_path.read_text(encoding="utf-8", errors="replace")
            for line in url_content.splitlines():
                if line.upper().startswith("URL="):
                    url = line[4:].strip()
                    output += f"\n\n---\n## Source Link\n**[🔗 View Original Webpage]({url})**\n"
                    break
        except Exception:
            pass
    else:
        # For all other files, append an embed link
        # Use ![[file]] for images/pdfs/md, and [[file]] for others like code
        if original_path.is_dir():
            output += f"\n\n---\n## Source Folder\n[[{original_path.name}]]\n"
        else:
            embed_prefix = "!" if original_path.suffix.lower() in [".png", ".jpg", ".jpeg", ".webp", ".pdf", ".mp3", ".mp4", ".wav", ".m4a", ".mov", ".md"] else ""
            output += f"\n\n---\n## Source File\n{embed_prefix}[[{original_path.name}]]\n"

    dest_file.write_text(output, encoding="utf-8")
    print(f"  ✅ Note → {wsl_to_windows(str(dest_file))}")

    # Process backlinks
    try:
        backlinks = []
        in_backlinks = False
        for line in output.splitlines():
            if line.startswith("update_backlinks:"):
                in_backlinks = True
            elif in_backlinks and line.startswith("  - "):
                backlinks.append(line[4:].strip().strip("'\"[]"))
            elif in_backlinks and not line.startswith(" "):
                in_backlinks = False

        if backlinks:
            cache = VaultIndexer(config["vault_root"]).get_cache()
            file_map = cache.get("files", {})
            for bl in backlinks:
                if bl in file_map:
                    target_path = resolve_path(config["vault_root"]) / file_map[bl]
                    if target_path.exists():
                        bl_content = target_path.read_text(encoding="utf-8")
                        bl_content += f"\n\n---\n**Related:** [[{dest_file.stem}]]\n"
                        target_path.write_text(bl_content, encoding="utf-8")
                        print(f"  🔗 Added backlink to {bl}")
                        logging.info(f"  Added backlink to {bl}")
    except Exception as e:
        logging.error(f"  Failed to process backlinks: {e}")

    # 2. Always move the original file to the Obsidian attachments folder (except shortcuts)
    if not config.get("skip_attachment", False) and original_path.suffix.lower() not in [".url", ".webloc"]:
        attach_dir = resolve_path(config["vault_root"]) / ATTACHMENTS_DIR
        attach_dir.mkdir(parents=True, exist_ok=True)
        attach_file = attach_dir / original_path.name
        
        if attach_file.exists():
            if original_path.is_dir():
                attach_file = attach_dir / f"{original_path.name}-{int(time.time())}"
            else:
                attach_file = attach_dir / f"{original_path.stem}-{int(time.time())}{original_path.suffix}"
            
        if original_path.is_dir():
            shutil.copytree(str(original_path), str(attach_file))
        else:
            shutil.copy2(original_path, attach_file)
        print(f"  ✅ Attachment → {wsl_to_windows(str(attach_file))}")
        
    sys.stdout.flush()

def move_to_failed(file_path: Path, config: dict):
    failed_dir = resolve_path(config["failed_dir"])
    failed_dir.mkdir(parents=True, exist_ok=True)
    dest = failed_dir / file_path.name
    # Append timestamp if it exists to avoid overwriting failed files
    if dest.exists():
        if file_path.is_dir():
            dest = failed_dir / f"{file_path.name}-{int(time.time())}"
        else:
            dest = failed_dir / f"{file_path.stem}-{int(time.time())}{file_path.suffix}"
    shutil.move(str(file_path), str(dest))
    logging.info(f"  ❌ Moved to Failed: {file_path.name}")

def archive_original(file_path: Path, config: dict):
    archive_dir = resolve_path(config["archive_dir"])
    archive_dir.mkdir(parents=True, exist_ok=True)
    dest = archive_dir / file_path.name
    if dest.exists():
        if file_path.is_dir():
            dest = archive_dir / f"{file_path.name}-{int(time.time())}"
        else:
            dest = archive_dir / f"{file_path.stem}-{int(time.time())}{file_path.suffix}"
    shutil.move(str(file_path), str(dest))
    logging.info(f"  📦 Archived: {file_path.name}")

# --- Event Handler -----------------------------------------------------------
class GeminiHandler(FileSystemEventHandler):
    def __init__(self, config):
        self.config = config
        self.processing = set()

    def on_created(self, event):
        self.process(Path(event.src_path))

    def process_directory(self, dir_path: Path):
        self.processing.add(dir_path)
        logging.info(f"DETECTED FOLDER: {dir_path.name}")
        print(f"\n📁 {dir_path.name}")
        sys.stdout.flush()
        
        # Wait for folder size to stabilize
        last_size = -1
        retries = 0
        while retries < 15:
            current_size = sum(f.stat().st_size for f in dir_path.rglob('*') if f.is_file())
            if current_size == last_size and current_size >= 0:
                if current_size > 0 or retries > 2:
                    break
            last_size = current_size
            time.sleep(2)
            retries += 1
        
        try:
            all_files = []
            for root, _, files in os.walk(dir_path):
                for f in files:
                    if not f.startswith("."):
                        all_files.append(Path(root) / f)
                        
            if not all_files:
                logging.warning(f"SKIPPING: {dir_path.name} (empty folder)")
                return

            code_extensions = {"py", "cpp", "h", "js", "ts", "ino", "sh", "yaml", "xml", "sql", "html", "css", "json", "c", "java", "go", "rs", "php", "rb"}
            code_count = sum(1 for f in all_files if f.suffix.lower().lstrip(".") in code_extensions)
            is_code_majority = (code_count / len(all_files)) >= 0.5
            
            destination = "coding 1s & 0s" if is_code_majority else "Research"
            target_folder = f"{destination}/{dir_path.name}"
            
            folder_config = self.config.copy()
            folder_config["smart_routing"] = False
            folder_config["skip_attachment"] = True
            folder_config["sibling_files"] = [f.name for f in all_files]
            
            for f in all_files:
                ext = f.suffix.lower().lstrip(".")
                route = None
                for r_name, r_data in self.config["routes"].items():
                    if ext in r_data["extensions"]:
                        route = r_data.copy()
                        route["name"] = r_name
                        break
                        
                if not route:
                    route = {
                        "name": "unknown",
                        "extensions": [],
                        "instructions": "Analyze this file.",
                    }
                
                route["destination"] = target_folder
                
                try:
                    content = None
                    if route.get("name") not in ["images", "media"] and f.suffix.lower() not in [".pdf", ".mhtml", ".mht", ".docx", ".pptx", ".xlsx"]:
                        content = read_file_content(f)
                        if not content and f.suffix.lower() not in [".url", ".webloc"]:
                            logging.warning(f"SKIPPING: {f.name} (no content extracted)")
                            continue
                            
                    logging.info(f"PROCESSING FOLDER FILE: {f.name} via Gemini ({route.get('name', 'unknown')})")
                    print(f"  🤖 Processing {f.name} via Gemini...")
                    sys.stdout.flush()
                    
                    output = process_with_gemini(f, content, route, folder_config)
                    route_output(f, output, route, folder_config)
                except Exception as file_e:
                    logging.error(f"  ❌ Error processing {f.name}: {file_e}")
                    print(f"  ❌ Error processing {f.name}: {file_e}")
                    sys.stdout.flush()
            
            # Copy whole directory to attachments
            attach_dir = resolve_path(self.config["vault_root"]) / ATTACHMENTS_DIR
            attach_dir.mkdir(parents=True, exist_ok=True)
            attach_file = attach_dir / dir_path.name
            if attach_file.exists():
                attach_file = attach_dir / f"{dir_path.name}-{int(time.time())}"
            shutil.copytree(str(dir_path), str(attach_file))
            print(f"  ✅ Folder Attachment → {wsl_to_windows(str(attach_file))}")
            
            if self.config["archive_originals"]:
                archive_original(dir_path, self.config)
                
        except Exception as e:
            logging.error(f"ERROR: {dir_path.name}: {e}")
            print(f"  ❌ Error: {e}")
            sys.stdout.flush()
            try:
                move_to_failed(dir_path, self.config)
            except Exception as move_e:
                logging.error(f"  Failed to move {dir_path.name} to Failed folder: {move_e}")
        finally:
            if dir_path in self.processing:
                self.processing.remove(dir_path)

    def process(self, file_path: Path):
        archive_dir = resolve_path(self.config["archive_dir"])
        failed_dir = resolve_path(self.config["failed_dir"])

        # Completely ignore the archive and failed directories and their contents
        if file_path.is_relative_to(archive_dir) or file_path.is_relative_to(failed_dir):
            return

        if file_path in self.processing or not file_path.exists():
            return
            
        if file_path.is_dir():
            self.process_directory(file_path)
            return
        
        ext = file_path.suffix.lower().lstrip(".")
        route = None
        for r_name, r_data in self.config["routes"].items():
            if ext in r_data["extensions"]:
                route = r_data
                route["name"] = r_name
                break
        
        if not route: return

        self.processing.add(file_path)
        logging.info(f"DETECTED: {file_path.name}")
        print(f"\n📥 {file_path.name}")
        sys.stdout.flush()
        time.sleep(2) # Give Windows extra time to finish writing the file
        
        try:
            content = None
            # Do not extract text for natively handled types (images, docs, pdfs, media, mhtml)
            if route["name"] not in ["images", "media"] and file_path.suffix.lower() not in [".pdf", ".mhtml", ".mht", ".docx", ".pptx", ".xlsx"]:
                content = read_file_content(file_path)
                if not content:
                    logging.warning(f"SKIPPING: {file_path.name} (no content extracted)")
                    return

            logging.info(f"PROCESSING: {file_path.name} via Gemini ({route['name']})")
            print(f"  🤖 Processing via Gemini ({route['name']})...")
            sys.stdout.flush()
            
            output = process_with_gemini(file_path, content, route, self.config)
            route_output(file_path, output, route, self.config)
            
            if self.config["archive_originals"]:
                archive_original(file_path, self.config)
                
        except Exception as e:
            logging.error(f"ERROR: {file_path.name}: {e}")
            print(f"  ❌ Error: {e}")
            sys.stdout.flush()
            # Move the file to the 'Failed' folder
            try:
                move_to_failed(file_path, self.config)
            except Exception as move_e:
                logging.error(f"  Failed to move {file_path.name} to Failed folder: {move_e}")
        finally:
            if file_path in self.processing:
                self.processing.remove(file_path)

# --- Main --------------------------------------------------------------------
def main():
    config = load_config()
    
    log_path = Path(__file__).parent / "watcher.log"
    logging.basicConfig(
        level=logging.INFO,
        format="%(asctime)s %(message)s",
        handlers=[
            logging.FileHandler(log_path),
            logging.StreamHandler(sys.stdout)
        ],
        force=True
    )
    
    watch_dir = resolve_path(config["watch_dir"])
    watch_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"👁  Gemini Vault Watcher Active")
    logging.info(f"STARTED: Monitoring {wsl_to_windows(str(watch_dir))}")
    print(f"   Watching: {wsl_to_windows(str(watch_dir))}")
    print(f"   Vault:    {wsl_to_windows(config['vault_root'])}")
    
    handler = GeminiHandler(config)
    observer = Observer()
    observer.schedule(handler, str(watch_dir), recursive=False)
    observer.start()

    print("   Status: Monitoring active (with 10s poll fallback)")
    sys.stdout.flush()

    try:
        while True:
            for item in watch_dir.iterdir():
                if not item.name.startswith("."):
                    handler.process(item)
            time.sleep(10)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()


if __name__ == "__main__":
    main()
